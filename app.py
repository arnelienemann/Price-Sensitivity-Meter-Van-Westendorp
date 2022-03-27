from email import header
from typing import final
import streamlit as st
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from functools import reduce
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor

st.header("Price Sensitivity Meter")

def pdf(df, col):
    
    stats_df = df.groupby(col)[col].agg('count').pipe(pd.DataFrame).rename(columns = {col: f'{col}_frequency'})
    stats_df[f'{col}_pdf'] = stats_df[f'{col}_frequency'] / sum(stats_df[f'{col}_frequency'])
    return stats_df

def cdf(df, col):

    # Frequency
    stats_df = df.groupby(col)[col].agg('count').pipe(pd.DataFrame).rename(columns = {col: f'{col}_frequency'})

    # PDF
    stats_df[f'{col}_pdf'] = stats_df[f'{col}_frequency'] / sum(stats_df[f'{col}_frequency'])

    # CDF
    stats_df[f'{col}_cdf'] = stats_df[f'{col}_pdf'].cumsum()
    stats_df.reset_index(inplace=True)
    stats_df.drop([f'{col}_frequency', f'{col}_pdf'], axis=1, inplace=True)
    stats_df.rename(columns = {col: 'Price', f'{col}_cdf': col}, inplace=True)
    
    return stats_df

def cdf_table(df, interpolate=False):
    '''
    Re-creating R's function output$data_vanwestendorp
    '''
    cdfs = [cdf(df, 'Too Cheap'), cdf(df, 'Cheap'), cdf(df, 'Expensive'), cdf(df, 'Too Expensive')]
    cdfs = reduce(lambda left, right: pd.merge(left, right, on=['Price'], how='outer'), cdfs).sort_values('Price')
    cdfs = cdfs.fillna(method='ffill').fillna(0)
    cdfs['Too Cheap'] = 1 - cdfs['Too Cheap']
    cdfs['Cheap'] = 1 - cdfs['Cheap']
    cdfs['Not Cheap'] = 1 - cdfs['Cheap']
    cdfs['Not Expensive'] = 1 - cdfs['Expensive']
    cdfs = cdfs.clip(lower=0)
    if interpolate == True:
        low = cdfs.Price.min()
        high = cdfs.Price.max()
        cdfs = pd.merge(pd.DataFrame(list(np.arange(low,high,0.01)), columns = ['Price']), cdfs, how='outer').sort_values('Price')
        cdfs['Price'] = cdfs['Price'].apply(lambda value: round(float(value),2))
        cdfs.drop_duplicates(['Price'], keep='last', inplace=True)
        cdfs = cdfs.interpolate(method ='linear', limit_direction ='forward')
        cdfs['Too Cheap'] = cdfs['Too Cheap'].fillna(1)
        cdfs['Cheap'] = cdfs['Cheap'].fillna(0)
        cdfs['Expensive'] = cdfs['Expensive'].fillna(0)
        cdfs['Too Expensive'] = cdfs['Too Expensive'].fillna(0)
        cdfs['Not Cheap'] = cdfs['Not Cheap'].fillna(0)
        cdfs['Not Expensive'] = cdfs['Not Expensive'].fillna(1)
        cdfs.reset_index(inplace=True)
        cdfs.drop('index', axis=1, inplace=True)
    return cdfs

data=pd.DataFrame()

uploaded_file = st.file_uploader("Choose a file")
if uploaded_file is not None:
    dataframe = pd.read_excel(uploaded_file, names=["Cheap","Expensive","Too Expensive", "Too Cheap"])
    data = dataframe

    #example data
    #data = pd.read_excel("data.xls", names=["Cheap","Expensive","Too Expensive", "Too Cheap"])
    data

    cdfs = cdf_table(data)

    Point_of_Marginal_Cheapness = cdfs.iloc[np.argwhere(np.diff(np.sign(cdfs['Too Cheap'] - cdfs['Not Cheap']))).flatten()+1]['Price'].values[0]
    Point_of_Marginal_Expensiveness = cdfs.iloc[np.argwhere(np.diff(np.sign(cdfs['Too Expensive'] - cdfs['Not Expensive']))).flatten()+1]['Price'].values[0]
    Optimal_Price_Point = cdfs.iloc[np.argwhere(np.diff(np.sign(cdfs['Too Expensive'] - cdfs['Too Cheap']))).flatten()+1]['Price'].values[0]

    st.write(f"Marginal Price Range: ${Point_of_Marginal_Cheapness:.2f} to ${Point_of_Marginal_Expensiveness:.2f}")
    st.write(f"Optimal Price Point: ${Optimal_Price_Point:.2f}")

    cdfs = cdfs.set_index("Price")
    #st.dataframe(cdfs[["Too Cheap", "Cheap", "Too Expensive", "Expensive"]]) # Drop cdfs["Not expensive"] & cdfs["Not Cheap"]
    st.write("\nPSM Results:")
    st.line_chart(cdfs[["Too Cheap", "Cheap", "Too Expensive", "Expensive"]])

    acceptance_data = pd.DataFrame()
    acceptance_data['Acceptance'] = 1 - cdfs['Cheap'] - cdfs['Expensive']
    acceptance_data['Revenue'] = acceptance_data['Acceptance'] * acceptance_data.index
    st.write("\nPrice acceptance:")
    st.line_chart(acceptance_data['Acceptance'])
    st.write("\nRevenue:")
    st.line_chart(acceptance_data['Revenue'])

if st.button('Create ppt'):

    file_path = "template - unbranded.pptx"
    prs = Presentation(file_path)

    title_slide_layout = prs.slide_layouts[0]
    divider_slide_layout = prs.slide_layouts[3]
    content_slide_layout = prs.slide_layouts[2]
    text_slide_layout = prs.slide_layouts[5]

    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "Price Sensitivity Meter"  + "\n" + "Survey Results"
    title_slide.shapes[1].text = "Arne Lienemann \nConsumer Insights | Global Portfolio Management"

    divider_slide = prs.slides.add_slide(divider_slide_layout)
    divider_slide.shapes.title.text = "Results"

    PSM_slide = prs.slides.add_slide(content_slide_layout)
    PSM_slide.shapes.title.text = "Results"

    chart_data = CategoryChartData()
    chart_data.categories = cdfs.index
    chart_data.add_series('Too Cheap', cdfs['Too Cheap'])
    chart_data.add_series('Cheap', cdfs['Cheap'])
    chart_data.add_series('Too Expensive', cdfs['Too Expensive'])
    chart_data.add_series('Expensive', cdfs['Expensive'])
    graphic_frame = PSM_slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.72), Inches(2.09), Inches(5.94), Inches(4.65), chart_data)
    graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
    graphic_frame.chart.value_axis.tick_labels.number_format = "0%"
    graphic_frame.chart.value_axis.tick_labels.font.size = Pt(10.5)
    #graphic_frame.chart.has_legend = False

    Acceptance_slide = prs.slides.add_slide(content_slide_layout)
    Acceptance_slide.shapes.title.text = "Price Acceptance"

    chart_data = CategoryChartData()
    chart_data.categories = acceptance_data.index
    chart_data.add_series('', acceptance_data['Acceptance'])
    graphic_frame = Acceptance_slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.72), Inches(2.09), Inches(5.94), Inches(4.65), chart_data)
    graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
    graphic_frame.chart.value_axis.tick_labels.number_format = "0%"
    graphic_frame.chart.value_axis.tick_labels.font.size = Pt(10.5)
    graphic_frame.chart.has_legend = False

    Revenue_slide = prs.slides.add_slide(content_slide_layout)
    Revenue_slide.shapes.title.text = "Revenue"

    chart_data = CategoryChartData()
    chart_data.categories = acceptance_data.index
    chart_data.add_series('', acceptance_data['Revenue'])
    graphic_frame = Revenue_slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.72), Inches(2.09), Inches(5.94), Inches(4.65), chart_data)
    graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
    #graphic_frame.chart.value_axis.tick_labels.number_format = "0%"
    graphic_frame.chart.value_axis.tick_labels.font.size = Pt(10.5)
    graphic_frame.chart.has_legend = False

    prs.save('temp_results.pptx')

    #with open(file_path, 'rb') as my_file:
    with open('temp_results.pptx', 'rb') as my_file:
        st.download_button(label = 'Download', data = my_file, file_name = 'Price Sensitivity Meter.pptx') 

#python -m streamlit run app.py